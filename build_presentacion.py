"""
Genera presentacion_tp3.pptx — segundo approach con diseño visual mejorado.
Cada slide tiene:
  · header bar superior con sección + número de slide
  · título + subtítulo
  · contenido principal en cards (no texto plano)
  · conclusion box inferior con el takeaway clave
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path("/Users/katiamenshikoff/ITBA/SIA/SIA-TP3")
OUT = ROOT / "presentacion_tp3.pptx"

# ===== Colors =====
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_TITLE = RGBColor(0x1F, 0x2A, 0x44)
C_TEXT = RGBColor(0x33, 0x33, 0x33)
C_MUTED = RGBColor(0x66, 0x6E, 0x80)
C_BORDER = RGBColor(0xE2, 0xE6, 0xEE)
C_CARD_BG = RGBColor(0xFA, 0xFB, 0xFD)
C_HP_BG = RGBColor(0xEE, 0xF3, 0xFA)
C_HP_BORDER = RGBColor(0xBF, 0xD2, 0xEC)
C_TABLE_HDR = RGBColor(0x1F, 0x2A, 0x44)

# Section colors
C_EJ1 = RGBColor(0x2E, 0x6F, 0xDB)
C_EJ2 = RGBColor(0x1F, 0x8E, 0x6F)
C_EJ3 = RGBColor(0xC4, 0x6A, 0x2E)
C_GEN = RGBColor(0x4A, 0x4E, 0x69)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ============================ HELPERS ============================

def text_set(tf, text, size=14, bold=False, color=C_TEXT, align=None, font=None):
    tf.word_wrap = True
    tf.clear()
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        if font:
            run.font.name = font


def header_bar(slide, section, page, color):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.4))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.35), Inches(0.06), Inches(10), Inches(0.28))
    text_set(t.text_frame, section, size=11, bold=True, color=WHITE)
    t = slide.shapes.add_textbox(Inches(11.4), Inches(0.06), Inches(1.6), Inches(0.28))
    text_set(t.text_frame, f"{page:02d}", size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def title_block(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.55))
    text_set(box.text_frame, title, size=22, bold=True, color=C_TITLE)
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.4), Inches(1.08), Inches(12.5), Inches(0.32))
        text_set(sb.text_frame, subtitle, size=13, color=C_MUTED)


def takeaway_box(slide, text, color, top=Inches(6.55)):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), top, Inches(12.5), Inches(0.78))
    box.adjustments[0] = 0.18
    box.fill.solid(); box.fill.fore_color.rgb = color; box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.28); tf.margin_right = Inches(0.28)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "CONCLUSIÓN   "
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = RGBColor(0xC8, 0xD8, 0xF0)
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(13); r2.font.bold = True; r2.font.color.rgb = WHITE


def hp_box(slide, lines, left=Inches(0.4), top=Inches(1.55), width=Inches(4.2), height=Inches(3.0)):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.05
    box.fill.solid(); box.fill.fore_color.rgb = C_HP_BG
    box.line.color.rgb = C_HP_BORDER; box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14); tf.margin_bottom = Inches(0.1)
    tf.word_wrap = True; tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "⚙   HIPERPARÁMETROS"
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = C_TITLE
    sp = tf.add_paragraph(); sp.add_run().text = " "
    sp.runs[0].font.size = Pt(4)
    for ln in lines:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = "•  " + ln
        r.font.size = Pt(11); r.font.color.rgb = C_TEXT


def insight_card(slide, title, body, left, top, width, height, accent_color=C_EJ1, body_size=11):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bg.fill.solid(); bg.fill.fore_color.rgb = C_CARD_BG
    bg.line.color.rgb = C_BORDER; bg.line.width = Pt(0.5)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    accent.fill.solid(); accent.fill.fore_color.rgb = accent_color; accent.line.fill.background()
    th = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.06), width - Inches(0.3), Inches(0.3))
    text_set(th.text_frame, title, size=11, bold=True, color=accent_color)
    tb = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.38), width - Inches(0.3), height - Inches(0.45))
    tf = tb.text_frame; tf.word_wrap = True; tf.clear()
    lines = body if isinstance(body, list) else body.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln
        r.font.size = Pt(body_size); r.font.color.rgb = C_TEXT


def stat_card(slide, value, label, left, top, width, height, color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.08
    box.fill.solid(); box.fill.fore_color.rgb = color; box.line.fill.background()
    tf = box.text_frame
    tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.06)
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = value
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run(); r.text = label
    r.font.size = Pt(10); r.font.color.rgb = RGBColor(0xE8, 0xEE, 0xF8)


def add_image(slide, img_path, left, top, width=None, height=None):
    p = Path(img_path)
    if not p.exists():
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width or Inches(6), height or Inches(4))
        ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        ph.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        text_set(ph.text_frame, f"[falta imagen]\n{p.name}", size=10, color=C_MUTED, align=PP_ALIGN.CENTER)
        return
    if width and height:
        slide.shapes.add_picture(str(p), left, top, width=width, height=height)
    elif width:
        slide.shapes.add_picture(str(p), left, top, width=width)
    elif height:
        slide.shapes.add_picture(str(p), left, top, height=height)
    else:
        slide.shapes.add_picture(str(p), left, top)


def add_table(slide, headers, rows, left, top, width, height, header_color=C_TABLE_HDR):
    tbl_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    tbl = tbl_shape.table
    # header row
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_color
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
        text_set(cell.text_frame, h, size=10, bold=True, color=WHITE)
    # body rows
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 1 else RGBColor(0xF6, 0xF8, 0xFC)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            text_set(cell.text_frame, str(val), size=10, color=C_TEXT)


def section_divider(slide, title, subtitle, color, page):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = color; bg.line.fill.background()
    # decorative thin top stripe
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.4), SLIDE_W, Inches(0.06))
    strip.fill.solid(); strip.fill.fore_color.rgb = WHITE; strip.line.fill.background()
    # eyebrow
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(0.5))
    text_set(t.text_frame, "TP3 · SIA 2026", size=13, bold=True, color=RGBColor(0xCC, 0xD8, 0xEE))
    # title
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.85), Inches(11.7), Inches(1.5))
    text_set(t.text_frame, title, size=44, bold=True, color=WHITE)
    # subtitle
    t = slide.shapes.add_textbox(Inches(0.8), Inches(4.45), Inches(11.7), Inches(1))
    text_set(t.text_frame, subtitle, size=18, color=RGBColor(0xEE, 0xF3, 0xFA))
    # page corner
    t = slide.shapes.add_textbox(Inches(11.4), Inches(6.95), Inches(1.6), Inches(0.35))
    text_set(t.text_frame, f"{page:02d}", size=12, bold=True, color=RGBColor(0xCC, 0xD8, 0xEE), align=PP_ALIGN.RIGHT)


# ============================ BUILD ============================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

PAGE = 0


def new_slide(section=None, color=None):
    """Create a slide and (optionally) draw the header chrome."""
    global PAGE
    PAGE += 1
    s = prs.slides.add_slide(BLANK)
    if section is not None and color is not None:
        header_bar(s, section, PAGE, color)
    return s


# ---------- PORTADA ----------
PAGE += 1
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid(); bg.fill.fore_color.rgb = C_TITLE; bg.line.fill.background()
# decorative blocks
for i, c in enumerate([C_EJ1, C_EJ2, C_EJ3]):
    bk = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4 + i * 0.4), Inches(0.4), Inches(0.3), Inches(0.18))
    bk.fill.solid(); bk.fill.fore_color.rgb = c; bk.line.fill.background()
t = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.5))
text_set(t.text_frame, "TP3 · ITBA · SIA 2026", size=14, bold=True, color=RGBColor(0xCC, 0xD8, 0xEE))
t = s.shapes.add_textbox(Inches(0.8), Inches(2.95), Inches(11.7), Inches(1.5))
text_set(t.text_frame, "Perceptrón Simple y Multicapa", size=42, bold=True, color=WHITE)
t = s.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.6))
text_set(t.text_frame, "Detección de fraude · Clasificación de dígitos · Hacia ≥ 98%", size=20, color=RGBColor(0xCC, 0xD8, 0xEE))
t = s.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4))
text_set(t.text_frame, "Segundo intento · Decisiones ancladas en las clases del curso", size=12, color=RGBColor(0xAA, 0xB8, 0xD0))

# ---------- AGENDA ----------
s = new_slide("AGENDA", C_GEN)
title_block(s, "Lo que vamos a recorrer", "Tres ejercicios · misma metodología · cadena de decisiones explícita")
# 3 cards horizontales
insight_card(s, "EJERCICIO 1 — Fraude",
    ["Distillation BigModel → TinyModel",
     "Perceptrón lineal vs no-lineal",
     "• Análisis del dataset",
     "• Sweep LR + threshold",
     "• Generalización"],
    Inches(0.4), Inches(1.7), Inches(4.2), Inches(4.6), C_EJ1, body_size=12)
insight_card(s, "EJERCICIO 2 — Dígitos (MLP)",
    ["Clasificación 10 clases · imágenes 28×28",
     "Exploración de LR · arch · optimizador",
     "• Análisis del dataset",
     "• 1ª tanda: sweeps one-at-a-time",
     "• 2ª tanda: experimentos cruzados",
     "• Convergencia + generalización"],
    Inches(4.7), Inches(1.7), Inches(4.2), Inches(4.6), C_EJ2, body_size=12)
insight_card(s, "EJERCICIO 3 — ≥ 98%",
    ["Subir test_acc sobre digits_test.csv",
     "Dos palancas:",
     "• Sumar more_digits.csv",
     "• Grid de regularización (L2 × σ)",
     "Comparativa final"],
    Inches(9.0), Inches(1.7), Inches(3.9), Inches(4.6), C_EJ3, body_size=12)
takeaway_box(s, "Cada decisión técnica está anclada en las clases del curso (regularización · métricas · optimizadores).", C_GEN)

# ===============================================================
# EJ1 — DIVISOR
# ===============================================================
PAGE += 1
s = prs.slides.add_slide(BLANK)
section_divider(s, "Ejercicio 1", "Detección de fraude — distillation BigModel → TinyModel · 7500 muestras · 11.59% fraude", C_EJ1, PAGE)

# ----- Ej1: análisis dataset -----
s = new_slide("EJ1 · Análisis del dataset", C_EJ1)
title_block(s, "Qué encontramos en `fraud_dataset.csv`", "Análisis univariado feature-vs-target antes de modelar")
add_image(s, ROOT/"ejercicio1/analisis_dataset/scatter_grid.png", Inches(0.4), Inches(1.55), width=Inches(7.0))
insight_card(s, "TARGET DE ENTRENAMIENTO",
    ["• big_model_fraud_probability (continuo [0,1])",
     "• Loss: MSE",
     "• Ground truth de evaluación: flagged_fraud"],
    Inches(7.7), Inches(1.55), Inches(5.3), Inches(1.5), C_EJ1)
insight_card(s, "3 REGLAS DURAS DETECTADAS",
    ["• amount_usd: mayoría de fraudes < 750",
     "• quantity_purchased ≥ 10 → casi todo fraude",
     "• session_duration_seconds < 150 → enriquecido"],
    Inches(7.7), Inches(3.15), Inches(5.3), Inches(1.5), C_EJ1)
insight_card(s, "BASELINE “3 REGLAS” (sin modelo)",
    ["Precision = 100% · Recall = 80%",
     "Accuracy = 97.68% · F1 ≈ 0.889",
     "Es el techo que tenemos que igualar."],
    Inches(7.7), Inches(4.75), Inches(5.3), Inches(1.6), C_EJ1)
takeaway_box(s, "Hay un baseline de 3 reglas duras con F1 ≈ 0.889 — el perceptrón sólo vale la pena si lo iguala.", C_EJ1)

# ----- Ej1: escalones -----
s = new_slide("EJ1 · Estructura del dataset", C_EJ1)
title_block(s, "Las 3 reglas son discontinuidades (escalones)", "Predicción teórica: el lineal va a underfittear estructuralmente")
add_image(s, ROOT/"ejercicio1/analisis_dataset/escalones.png", Inches(0.4), Inches(1.55), width=Inches(7.5))
insight_card(s, "HIPÓTESIS TEÓRICA",
    ["Saltos de ~6% → 100% en un valor entero.",
     "Geométricamente imposibles para una",
     "recta + sigmoide saturada.",
     "",
     "• Lineal → underfitting estructural",
     "• No-lineal (sigmoide) → puede acercarse"],
    Inches(8.1), Inches(1.55), Inches(4.9), Inches(2.5), C_EJ1)
insight_card(s, "DEFINICIÓN DE UNDERFITTING",
    ["No alcanzar el baseline en train:",
     "  acc ≥ 97.68%   ·   prec ≥ 100%   ·   rec ≥ 80%",
     "",
     "Como mínimo queremos igualar el RECALL",
     "para que valga la pena usar el perceptrón."],
    Inches(8.1), Inches(4.1), Inches(4.9), Inches(2.25), C_EJ1)
takeaway_box(s, "Las 3 reglas son escalones → predecimos que el lineal va a quedarse corto en MSE; el no-lineal debería acercarse.", C_EJ1)

# ----- Ej1: pipeline / decisiones -----
s = new_slide("EJ1 · Pipeline y decisiones", C_EJ1)
title_block(s, "Pipeline y decisiones de diseño", "Todo anclado en las clases del curso")
insight_card(s, "PREPROCESAMIENTO",
    ["• Z-score fit-on-train por fold",
     "  (sin data leakage)",
     "• Bipolar no aplica (problema binario",
     "  con probabilidades del BigModel)"],
    Inches(0.4), Inches(1.55), Inches(4.0), Inches(2.2), C_EJ1)
insight_card(s, "VALIDACIÓN",
    ["• K-fold estratificado, K = 5",
     "  (mantiene el 11.59% de positivos)",
     "• Multi-seed × 5 → separa varianza por",
     "  inicialización vs por partición",
     "• Total: 25 corridas por config"],
    Inches(4.55), Inches(1.55), Inches(4.0), Inches(2.2), C_EJ1)
insight_card(s, "MODELOS",
    ["• Lineal (ADALINE, identidad)",
     "  entrenado online con MSE",
     "• No-lineal (sigmoide), mismo MSE,",
     "  misma data — sólo cambia la activación",
     "• Init: U(−0.1, 0.1), seed + k_fold"],
    Inches(8.7), Inches(1.55), Inches(4.3), Inches(2.2), C_EJ1)
insight_card(s, "CRITERIO DE CORTE",
    ["• Por techo de épocas — epsilon NUNCA dispara",
     "  (calibrado por debajo del MSE asintótico).",
     "• Convergencia honesta = plateau empírico:",
     "  pendiente del MSE en las últimas 50 ép. ≈ 0."],
    Inches(0.4), Inches(3.9), Inches(6.2), Inches(2.4), C_EJ1)
insight_card(s, "THRESHOLD (post-training)",
    ["• Vive POST-training: no cambia los pesos,",
     "  sólo binariza la salida continua.",
     "• Sweep denso de thr ∈ [0, 1].",
     "• thr* = el que maximiza F1 sobre las 25 corridas."],
    Inches(6.7), Inches(3.9), Inches(6.3), Inches(2.4), C_EJ1)
takeaway_box(s, "K=5 estratificado × 5 seeds + z-score fit-on-train + threshold post-training → 25 estimaciones por config sin leakage.", C_EJ1)

# ----- Ej1: K-fold sweep -----
s = new_slide("EJ1 · Elección de K", C_EJ1)
title_block(s, "¿Por qué K = 5?", "Sweep empírico K ∈ {2, 3, 5, 10} sobre el lineal")
hp_box(s, [
    "Modelo: perceptrón lineal",
    "LR = 1e-4 · epochs = 500",
    "Threshold = 0.69 (thr* del LR)",
    "Seed = 42 (única, comparable fold-a-fold)",
    "Estratificado por flagged_fraud",
], top=Inches(1.55), height=Inches(2.4))
add_image(s, ROOT/"Notas/ejercicio 1/Experimentos y analisis/LINEAR perceptron/kfold_sweep/kfold_sweep.png",
          Inches(4.9), Inches(1.55), width=Inches(8.0))
insight_card(s, "TRADE-OFF SESGO ↔ VARIANZA DEL ESTIMADOR",
    ["• K chico: train pequeño → estimación pesimista.",
     "• K grande: folds de test chicos → pocas muestras de fraude → métricas ruidosas.",
     "• K=5 minimiza std(MSE) = 0.00076; K=10 empeora a 0.00206 sin mejorar la media."],
    Inches(0.4), Inches(4.1), Inches(12.6), Inches(2.2), C_EJ1)
takeaway_box(s, "K=5 minimiza la varianza del estimador para 7500 filas con 11.59% positivos. K=10 duplica cómputo sin ganancia.", C_EJ1)

# ----- Ej1: sweep LR lineal — convergencia -----
s = new_slide("EJ1 · Sweep LR · lineal · convergencia", C_EJ1)
title_block(s, "Sweep LR — perceptrón lineal · convergencia", "5 seeds × 3 LRs × 5 folds = 75 corridas")
hp_box(s, [
    "Modelo: lineal (identidad, MSE)",
    "LRs: {1e-5, 1e-4, 1e-3}",
    "Seeds: 7, 13, 21, 42, 99",
    "K-fold = 5 estratificado",
    "Epochs = 500 · entrenamiento online",
    "Init: U(−0.1, 0.1) · z-score fit-on-train",
], top=Inches(1.55), height=Inches(3.0))
add_image(s, ROOT/"Notas/ejercicio 1/Experimentos y analisis/LINEAR perceptron/sweep_lr/multiseed/convergence.png",
          Inches(4.9), Inches(1.55), width=Inches(8.0))
insight_card(s, "ARGUMENTO DE CONVERGENCIA",
    ["75 corridas llegaron a las 500 ép. — epsilon NUNCA disparó.",
     "Tail-slope (pendiente últimas 50 ép.) ≈ 1e-13 → ruido numérico float64 → ΔMSE/Δep → 0 → PLATEAU."],
    Inches(0.4), Inches(4.7), Inches(12.6), Inches(1.6), C_EJ1)
takeaby = "Convergencia honesta = plateau (no epsilon). Argumento sólido para defender por qué decimos que el modelo “converge”."
takeaway_box(s, takeaby, C_EJ1)

# ----- Ej1: sweep LR lineal — métricas y threshold -----
s = new_slide("EJ1 · Sweep LR · lineal · métricas a thr*", C_EJ1)
title_block(s, "Métricas del sweep LR lineal", "Cada LR evaluado a SU threshold óptimo")
add_image(s, ROOT/"Notas/ejercicio 1/Experimentos y analisis/LINEAR perceptron/sweep_lr/multiseed/threshold_curves.png",
          Inches(0.4), Inches(1.55), width=Inches(7.0))
hp_box(s, [
    "thr* por LR (maximiza F1):",
    "  LR=1e-5 → thr*=0.69",
    "  LR=1e-4 → thr*=0.69",
    "  LR=1e-3 → thr*=0.78",
    "Métricas: mean ± std sobre 25 corridas",
    "Loss reportada: MSE (la que se minimizó)",
], left=Inches(7.6), top=Inches(1.55), width=Inches(5.4), height=Inches(2.8))
add_table(s,
    ["LR", "thr*", "MSE test", "Acc", "Prec", "Rec", "F1"],
    [
        ["1e-5", "0.69", "0.0262 ± 0.0012", "0.9741 ± 0.003", "0.925 ± 0.018", "0.846 ± 0.029", "0.883 ± 0.016"],
        ["1e-4", "0.69", "0.0265 ± 0.0014", "0.9736 ± 0.003", "0.910 ± 0.021", "0.857 ± 0.028", "0.883 ± 0.016"],
        ["1e-3", "0.78", "0.0460 ± 0.0062", "0.9713 ± 0.004", "0.905 ± 0.019", "0.841 ± 0.027", "0.872 ± 0.017"],
    ],
    Inches(0.4), Inches(4.6), Inches(12.6), Inches(1.7))
takeaway_box(s, "LR = 1e-4 (converge antes que 1e-5 con métricas equivalentes). LR = 1e-3 sufre SUBAJUSTE en MSE.", C_EJ1)

# ----- Ej1: sweep LR no-lineal — convergencia -----
s = new_slide("EJ1 · Sweep LR · no-lineal · convergencia", C_EJ1)
title_block(s, "Sweep LR — perceptrón no-lineal · convergencia", "Mismo protocolo · salida sigmoide")
hp_box(s, [
    "Modelo: no-lineal (sigmoide, MSE)",
    "LRs: {1e-3, 1e-2, 1e-1}",
    "Seeds: 7, 13, 21, 42, 99",
    "K-fold = 5 estratificado",
    "Epochs = 500 · entrenamiento online",
    "Init: U(−0.1, 0.1) · z-score fit-on-train",
], top=Inches(1.55), height=Inches(3.0))
add_image(s, ROOT/"Notas/ejercicio 1/Experimentos y analisis/NON LINEAR perceptron/sweep_lr/multiseed/convergence.png",
          Inches(4.9), Inches(1.55), width=Inches(8.0))
stat_card(s, "0.011", "MSE no-lineal", Inches(0.4), Inches(4.7), Inches(3.0), Inches(1.5), C_EJ1)
stat_card(s, "0.026", "MSE lineal",    Inches(3.6), Inches(4.7), Inches(3.0), Inches(1.5), C_EJ1)
stat_card(s, "≈ 2×",  "ganancia del no-lineal", Inches(6.8), Inches(4.7), Inches(3.0), Inches(1.5), C_EJ1)
takeaway_box(s, "El no-lineal converge a un MSE 2× mejor — la sigmoide aproxima los escalones que la recta no puede capturar.", C_EJ1)

# ----- Ej1: threshold sweep no-lineal -----
s = new_slide("EJ1 · Threshold sweep · no-lineal", C_EJ1)
title_block(s, "Threshold sweep no-lineal · curva PR", "El threshold es decisión de NEGOCIO, no HP técnico")
add_image(s, ROOT/"Notas/ejercicio 1/Experimentos y analisis/NON LINEAR perceptron/sweep_lr/multiseed/threshold_curves.png",
          Inches(0.4), Inches(1.55), width=Inches(6.3))
add_image(s, ROOT/"Notas/ejercicio 1/Experimentos y analisis/NON LINEAR perceptron/sweep_lr/multiseed/pr_curve.png",
          Inches(6.85), Inches(1.55), width=Inches(6.1))
insight_card(s, "CRITERIO (declarado antes de mirar)",
    ["Queremos máxima recall posible sin perder precisión. Con thr bajos, recall ≈ 1 pero precision colapsa (flag a casi todo).",
     "Usamos F1 como métrica síntesis para decidir thr* que balancea ambos costos.   →   thr* del LR=1e-2 ganador = 0.89."],
    Inches(0.4), Inches(5.3), Inches(12.6), Inches(1.1), C_EJ1)
takeaway_box(s, "thr* = 0.89 maximiza F1. El cliente puede mover el threshold según el costo relativo FN vs FP.", C_EJ1)

# ----- Ej1: comparativa lineal vs no-lineal vs baseline -----
s = new_slide("EJ1 · Comparativa final", C_EJ1)
title_block(s, "Lineal vs No-lineal vs Baseline (cada uno a su thr*)", "La métrica del entrenamiento (MSE) ≠ la métrica de decisión (F1)")
add_table(s,
    ["Modelo", "thr*", "MSE test", "Accuracy", "Precision", "Recall", "F1"],
    [
        ["Lineal (LR=1e-4)",     "0.69", "0.0265 ± 0.0014", "0.9736 ± 0.0035", "0.9100 ± 0.021", "0.8573 ± 0.028", "0.8825 ± 0.016"],
        ["No-lineal (LR=1e-2)",  "0.89", "0.0110 ± 0.0006", "0.9709 ± 0.0045", "0.8859 ± 0.021", "0.8594 ± 0.027", "0.8722 ± 0.020"],
        ["Baseline 3 reglas",    "—",    "—",              "0.9768",          "1.000",          "0.800",           "0.889"],
    ],
    Inches(0.4), Inches(1.55), Inches(12.6), Inches(2.2))
insight_card(s, "OBSERVACIÓN NO TRIVIAL",
    ["• El no-lineal gana en MSE por 2× (confirma la predicción teórica: aprende los escalones).",
     "• Pero en F1 quedan EMPATADOS, ambos rozando el baseline (0.889).",
     "• MSE evalúa cuán bien aproxima la probabilidad continua; F1 evalúa cuán bien clasifica binariamente.",
     "  Son objetos distintos y se contradicen — exactamente la regla 4 del CLAUDE.md."],
    Inches(0.4), Inches(4.0), Inches(12.6), Inches(2.3), C_EJ1)
takeaway_box(s, "Mejor MSE ≠ mejor F1. Ambos modelos rozan el baseline en F1 — la capacidad extra del no-lineal aprovecha en MSE, no en F1.", C_EJ1)

# ----- Ej1: generalización -----
s = new_slide("EJ1 · Generalización", C_EJ1)
title_block(s, "Generalización — sin overfitting", "Modelo elegido para CompanyX: no-lineal, LR=1e-2, thr*=0.89")
add_image(s, ROOT/"Notas/ejercicio 1/Generalizacion/gap_train_test.png", Inches(0.4), Inches(1.55), width=Inches(7.0))
hp_box(s, [
    "5 seeds × 5 folds = 25 corridas",
    "z-score fit-on-train por fold",
    "Epochs = 500 · sin early stopping",
    "thr* aplicado post-training",
    "Métricas a thr*=0.89:",
    "   MSE test = 0.01099 ± 0.00058",
    "   Acc = 0.9709 · F1 = 0.8722",
], left=Inches(7.7), top=Inches(1.55), width=Inches(5.3), height=Inches(3.0))
insight_card(s, "GAP TRAIN ↔ TEST",
    ["Lineal:    +0.00019    No-lineal: +0.00006",
     "Puntos sobre la diagonal y=x → MSE_train ≈ MSE_test → no hay memorización.",
     "El error que queda no es gap train/test sino UNDERFITTING — capacidad insuficiente, no exceso."],
    Inches(0.4), Inches(4.7), Inches(12.6), Inches(1.7), C_EJ1)
takeaway_box(s, "El modelo no memoriza: lo que aprendió en 6000 muestras aplica igual a las 1500 que no vio.", C_EJ1)

# ===============================================================
# EJ2 — DIVISOR
# ===============================================================
PAGE += 1
s = prs.slides.add_slide(BLANK)
section_divider(s, "Ejercicio 2", "Clasificación de dígitos · MLP · 12.449 imágenes 28×28 · LR · arch · optimizador", C_EJ2, PAGE)

# ----- Ej2: distribución -----
s = new_slide("EJ2 · Análisis del dataset", C_EJ2)
title_block(s, "Distribución de clases en digits.csv", "Hallazgos clave antes de modelar")
add_image(s, ROOT/"Notas/ejercicio 2/analisis_dataset/distribucion_clases.png", Inches(0.4), Inches(1.55), width=Inches(7.0))
insight_card(s, "DATASET",
    ["12.449 muestras · 784 features (28×28)",
     "Valores en [0, 1] · 9 clases presentes"],
    Inches(7.7), Inches(1.55), Inches(5.3), Inches(1.4), C_EJ2)
insight_card(s, "ANOMALÍA #1 — CLASE 8 AUSENTE",
    ["digits.csv no tiene NI UN ejemplo de la clase 8.",
     "→ capa de salida = 9 neuronas durante HP search.",
     "(Ej3 va a sumar more_digits.csv con 585 ejemplos)"],
    Inches(7.7), Inches(3.0), Inches(5.3), Inches(1.55), C_EJ2)
insight_card(s, "ANOMALÍA #2 — CLASE 5 DESBALANCEADA",
    ["271 ejemplos vs ~1500 las otras (6× menos).",
     "Implica: macro_F1 obligatorio (no sólo accuracy)."],
    Inches(7.7), Inches(4.65), Inches(5.3), Inches(1.6), C_EJ2)
takeaway_box(s, "Dos anomalías estructurales del dataset: clase 8 AUSENTE y clase 5 desbalanceada → marcan toda la metodología.", C_EJ2)

# ----- Ej2: muestras intra-clase -----
s = new_slide("EJ2 · Variabilidad intra-clase", C_EJ2)
title_block(s, "Muestras por clase", "Lo que el MLP tiene que aprender a clasificar")
add_image(s, ROOT/"Notas/ejercicio 2/analisis_dataset/muestras_por_clase.png", Inches(0.7), Inches(1.55), width=Inches(11.8))
takeaway_box(s, "Variabilidad intra-clase amplia (estilos, inclinaciones, grosor) — eso es lo que el MLP debe aprender a ignorar.", C_EJ2)

# ----- Ej2: media antes -----
s = new_slide("EJ2 · Imagen media · ANTES de normalizar", C_EJ2)
title_block(s, "Imagen media por clase — datos crudos", "Promedio píxel a píxel de las muestras de cada clase")
add_image(s, ROOT/"Notas/ejercicio 2/analisis_dataset/media_por_clase.png", Inches(0.4), Inches(1.55), width=Inches(8.4))
insight_card(s, "QUÉ MIDE",
    ["Si las muestras de una clase comparten",
     "forma típica → la media es nítida.",
     "Si hay alta variabilidad → media borrosa."],
    Inches(9.0), Inches(1.55), Inches(4.0), Inches(2.0), C_EJ2)
insight_card(s, "QUÉ VEMOS",
    ["Todas las medias son nítidas y",
     "reconocibles como el dígito correspondiente",
     "→ señal suficiente dentro de cada clase",
     "para que el MLP pueda aprender."],
    Inches(9.0), Inches(3.7), Inches(4.0), Inches(2.5), C_EJ2)
takeaway_box(s, "Estructura compartida dentro de cada clase → hay señal aprendible. El problema no es separabilidad, es generalizar a estilos nuevos.", C_EJ2)

# ----- Ej2: media después -----
s = new_slide("EJ2 · Imagen media · DESPUÉS de z-score", C_EJ2)
title_block(s, "Imagen media por clase — espacio normalizado", "x' = (x − μⱼ) / σⱼ · fit-on-train para evitar leakage")
add_image(s, ROOT/"Notas/ejercicio 2/analisis_dataset/media_normalizada.png", Inches(0.4), Inches(1.55), width=Inches(8.4))
insight_card(s, "CÓMO LEER LOS COLORES",
    ["• Rojo (z > 0): píxeles que esta clase",
     "  activa MÁS que el promedio.",
     "• Azul (z < 0): píxeles que activa MENOS.",
     "→ visualiza qué regiones distinguen cada clase."],
    Inches(9.0), Inches(1.55), Inches(4.0), Inches(2.4), C_EJ2)
insight_card(s, "QUÉ GANAMOS CON LA NORMALIZACIÓN",
    ["• Todos los píxeles en escala comparable",
     "  (media 0, std 1).",
     "• Estabiliza el entrenamiento por gradientes.",
     "• Permite que He/Xavier inicialicen como",
     "  asume su derivación teórica."],
    Inches(9.0), Inches(4.05), Inches(4.0), Inches(2.2), C_EJ2)
takeaway_box(s, "Z-score CONSERVA la estructura (formas reconocibles) y estabiliza el entrenamiento.", C_EJ2)

# ----- Ej2: decisiones del preánalisis -----
s = new_slide("EJ2 · Decisiones que surgen del preánalisis", C_EJ2)
title_block(s, "Cómo el dataset condiciona la metodología", "Cada decisión deriva de las anomalías encontradas")
add_table(s,
    ["Decisión", "Justificación"],
    [
        ["Capa de salida: 9 neuronas (HP search)",      "Sólo hay 9 clases en digits.csv"],
        ["Métrica principal: val_acc + macro_F1",        "Dataset desbalanceado (clase 5, n=271)"],
        ["Reportar Precision/Recall por clase",          "Detectar si el modelo falla en la clase 5"],
        ["Normalización z-score fit-on-train",           "Evita data leakage en K-fold"],
        ["K-fold estratificado",                          "Mantiene la proporción de la clase 5 en cada fold"],
        ["digits_test.csv intocado durante HP search",   "Es producción — se evalúa UNA sola vez al cerrar Ej2"],
    ],
    Inches(0.4), Inches(1.55), Inches(12.6), Inches(2.7))
insight_card(s, "POR QUÉ MACRO Y NO MICRO / WEIGHTED",
    ["• micro_f1 ≡ accuracy en multiclase con argmax → redundante.",
     "• weighted_f1 pondera por frecuencia → penaliza menos los errores en clase rara (LO OPUESTO de lo que queremos).",
     "• macro_f1 pesa igual cada clase → es el que defendemos (regla 4 del CLAUDE.md, clase de métricas)."],
    Inches(0.4), Inches(4.35), Inches(12.6), Inches(2.0), C_EJ2)
takeaway_box(s, "El dataset dicta la metodología — macro_F1 y estratificación no son “por defecto”, son consecuencias del desbalance.", C_EJ2)

# ----- Ej2: 1ra tanda — optimizers -----
s = new_slide("EJ2 · 1ª tanda · sweep de optimizadores", C_EJ2)
title_block(s, "Sweep de optimizadores", "Comparar SGD vs Momentum vs Adam con arch y batch fijos")
hp_box(s, [
    "Arch: [784, 100, 50, 10]",
    "Activaciones: relu, relu, softmax",
    "Loss: cross-entropy (+ softmax)",
    "LRs: 1e-4, 5e-4, 1e-3, 5e-3, 1e-2",
    "Batch = 32 · max epochs = 70",
    "Seeds × Folds = 5 × 5 = 25",
    "Opts: SGD · Momentum (β=0.9) ·",
    "       Adam (β1=0.9, β2=0.999)",
], top=Inches(1.55), height=Inches(3.4))
add_image(s, ROOT/"Notas/ejercicio 2/Primera tanda de experimentos/Optimizer/val_acc_vs_epoch_per_optimizer.png",
          Inches(4.9), Inches(1.55), width=Inches(8.0))
insight_card(s, "HALLAZGO QUE ABRE LA 2ª TANDA",
    ["Adam converge en 3-5 épocas vs ~20 de Momentum vs >40 de SGD.",
     "Pero Adam se desestabiliza con LR ≥ 5e-3 mientras SGD/Momentum siguen tolerando LR alto.",
     "→ Sospecha: el techo de Adam quizás no es del optimizador, sino del producto LR×batch (regla de la clase de optimizadores)."],
    Inches(0.4), Inches(5.05), Inches(12.6), Inches(1.4), C_EJ2)
takeaway_box(s, "Sospecha levantada: ¿el techo de Adam es del optimizador o del producto LR×batch? — motiva la 2ª tanda.", C_EJ2)

# ----- Ej2: Pre LR×Batch×Opt -----
s = new_slide("EJ2 · 2ª tanda · Pre LR × Batch × Opt", C_EJ2)
title_block(s, "Pre-experimento LR × Batch × Opt", "Confirmar empíricamente la regla LR ↔ batch")
hp_box(s, [
    "Stage 1 del cross_v1",
    "Arch: base [784,100,50,10]",
    "LRs: 5e-4, 1e-3, 5e-3",
    "Batches: 16, 32, 64, 128, 256",
    "Opts: SGD, Momentum, Adam",
    "Seeds × Folds = 3 × 5",
    "max epochs adaptado por (opt, LR)",
    "ES patience = 20",
], top=Inches(1.55), height=Inches(3.4))
add_image(s, ROOT/"Notas/ejercicio 2/Segunda tanda de experimentos/Pre_LR_Batch_Opt/stage1_heatmap_val_acc.png",
          Inches(4.9), Inches(1.55), width=Inches(8.0))
insight_card(s, "RESULTADO CLAVE",
    ["Adam@5e-3 con batch=256 es estable (val_loss=0.191). Con batch=16 no converge (val_loss=0.550).",
     "→ Confirma la sospecha: el techo de Adam NO es del optimizador, es del PRODUCTO LR×batch."],
    Inches(0.4), Inches(5.05), Inches(12.6), Inches(1.4), C_EJ2)
takeaway_box(s, "Confirmado: lo que rompía Adam@LR alto era LR×batch fuera de rango. Decisión: batch=64 para Adam@1e-3 en el grid principal.", C_EJ2)

# ----- Ej2: lr_batch_relationship -----
s = new_slide("EJ2 · Regla LR ↔ batch", C_EJ2)
title_block(s, "Regla LR ↔ batch — visualización", "Sólo Adam tiene sensibilidad fuerte al producto LR×batch")
add_image(s, ROOT/"Notas/ejercicio 2/Segunda tanda de experimentos/Pre_LR_Batch_Opt/lr_batch_relationship.png",
          Inches(0.4), Inches(1.55), width=Inches(8.3))
insight_card(s, "ADAM@5e-3 — DEPENDE DEL BATCH",
    ["batch 16 → 32 → 256",
     "val_loss: 0.55 → 0.36 → 0.19",
     "val_acc:  0.934 → 0.948 → 0.954"],
    Inches(9.0), Inches(1.55), Inches(4.0), Inches(2.0), C_EJ2)
insight_card(s, "ADAM EN SU ÓPTIMO",
    ["Adam@5e-4 y Adam@1e-3 ya estables.",
     "Curvas planas vs batch → no mueve la aguja",
     "porque el LR ya era el correcto."],
    Inches(9.0), Inches(3.7), Inches(4.0), Inches(1.9), C_EJ2)
insight_card(s, "SGD / MOMENTUM",
    ["Sensibilidad MUCHO menor al batch.",
     "Más robustos a errores en HP."],
    Inches(9.0), Inches(5.65), Inches(4.0), Inches(0.75), C_EJ2)
takeaway_box(s, "Patrón exactamente predicho por la clase de optimizadores: Adam adapta el paso por parámetro → es sensible a LR×batch.", C_EJ2)

# ----- Ej2: grid 3D -----
s = new_slide("EJ2 · Grid principal LR × Opt × Arch (stage 2)", C_EJ2)
title_block(s, "Grid 3D — LR × Opt × Arch", "60 celdas = 5 LRs × 3 opts × 4 archs · 3 seeds × k=5 = 15 corridas/celda")
hp_box(s, [
    "Archs:",
    "  shallow [784,128,10]",
    "  base    [784,100,50,10]",
    "  wider   [784,200,100,10]",
    "  deeper  [784,128,64,32,10]",
    "Activaciones: relu (+ softmax)",
    "Init: He (auto)",
    "Batch = best_batch[(opt,LR)] del stage 1",
    "ES patience = 20",
], top=Inches(1.55), height=Inches(3.8))
add_image(s, ROOT/"Notas/ejercicio 2/Segunda tanda de experimentos/Cross_LR_Opt_Arch/stage2_val_acc_vs_lr_per_opt.png",
          Inches(4.9), Inches(1.55), width=Inches(8.0))
insight_card(s, "VISTA MARGINAL POR OPTIMIZER (best-over-arch)",
    ["Adam pica en 1e-3 en las 4 archs (rango operativo estrecho).",
     "SGD pica en 1e-2 (plateau desde 5e-4 — robusto al HP).",
     "Momentum se desplaza según arch — única interacción fuerte del grid."],
    Inches(0.4), Inches(5.5), Inches(12.6), Inches(1.0), C_EJ2)
takeaway_box(s, "Adam domina pero es FRÁGIL: rango operativo estrecho. SGD/Momentum perdonan errores en HP pero techo más bajo.", C_EJ2)

# ----- Ej2: LR × OPT (4 paneles) -----
s = new_slide("EJ2 · LR × Opt (un panel por arch)", C_EJ2)
title_block(s, "LR × Opt — desagregado por arquitectura", "¿Hay interacción LR × Opt × Arch?")
add_image(s, ROOT/"Notas/ejercicio 2/Segunda tanda de experimentos/Cross_LR_Opt_Arch/lr_opt_val_acc_4panels.png",
          Inches(0.4), Inches(1.55), width=Inches(8.2))
insight_card(s, "SGD",
    ["Patrón idéntico en las 4 archs.",
     "Salto grande 1e-4 → 5e-4, después plateau.",
     "SIN interacción con arch."],
    Inches(8.8), Inches(1.55), Inches(4.2), Inches(1.5), C_EJ2)
insight_card(s, "MOMENTUM — interacción CLARA",
    ["shallow/wider (última hidden 128) toleran LR=1e-2.",
     "base/deeper (capa más estrecha) se rompen.",
     "Mecanismo: capa estrecha = cada peso pesa",
     "más en la salida → LR alto desestabiliza."],
    Inches(8.8), Inches(3.15), Inches(4.2), Inches(1.85), C_EJ2)
insight_card(s, "ADAM",
    ["Shapes idénticas en las 4 archs:",
     "pico en 1e-3, caída después.",
     "Lo que cambia es la MAGNITUD del pico."],
    Inches(8.8), Inches(5.1), Inches(4.2), Inches(1.3), C_EJ2)
takeaway_box(s, "Última capa hidden ancha actúa como “buffer” ante LR alto. Es el mecanismo estructural de la interacción LR×Opt×Arch.", C_EJ2)

# ----- Ej2: val_loss CE -----
s = new_slide("EJ2 · val_loss CE captura lo que la accuracy esconde", C_EJ2)
title_block(s, "val_loss CE vs val_acc", "Por qué reportamos las DOS (loss + métrica de clasificación)")
add_image(s, ROOT/"Notas/ejercicio 2/Segunda tanda de experimentos/Cross_LR_Opt_Arch/lr_opt_val_loss_4panels.png",
          Inches(0.4), Inches(1.55), width=Inches(8.2))
stat_card(s, "−1 pp",  "Δ val_acc (1e-3 → 1e-2)",    Inches(8.8), Inches(1.55), Inches(4.2), Inches(1.4), C_EJ2)
stat_card(s, "+46%",   "Δ val_loss CE (1e-3 → 1e-2)", Inches(8.8), Inches(3.1),  Inches(4.2), Inches(1.4), C_EJ2)
insight_card(s, "POR QUÉ PASA",
    ["El argmax sigue cayendo a veces en la clase correcta",
     "aunque las probabilidades estén totalmente desordenadas.",
     "CE detecta la descalibración antes que la accuracy."],
    Inches(8.8), Inches(4.65), Inches(4.2), Inches(1.75), C_EJ2)
takeaway_box(s, "CE responde “qué tan bien calibradas están las probabilidades”; accuracy responde “cuántas argmax acerté”. Son distintas (regla 4).", C_EJ2)

# ----- Ej2: Arch × Opt heatmap -----
s = new_slide("EJ2 · Arch × Opt (best-over-LR)", C_EJ2)
title_block(s, "¿Hay arquitecturas especialistas de un optimizador?", "Marginalizamos por “best-over-LR” (no “mean-over-LR”) para no contaminar con LRs sub-óptimos")
add_image(s, ROOT/"Notas/ejercicio 2/Segunda tanda de experimentos/Cross_LR_Opt_Arch/arch_opt_best_lr_heatmap.png",
          Inches(0.4), Inches(1.55), width=Inches(7.3))
insight_card(s, "1 · ADAM DOMINA LAS 4 ARCHS",
    ["Entre +0.003 y +0.004 sobre el 2° mejor opt.",
     "Ninguna arch prefiere Mom o SGD sobre Adam."],
    Inches(7.9), Inches(1.55), Inches(5.1), Inches(1.3), C_EJ2)
insight_card(s, "2 · WIDER ES ESPECIALISTA DE ADAM",
    ["Techo absoluto (0.9583) con Adam@1e-3.",
     "Pero queda 2°/3° con SGD y Momentum."],
    Inches(7.9), Inches(2.95), Inches(5.1), Inches(1.3), C_EJ2)
insight_card(s, "3 · SHALLOW ES “TODO TERRENO”",
    ["★ en SGD y ★ en Momentum.",
     "Empata con wider en Adam (tiebreaker más tarde)."],
    Inches(7.9), Inches(4.35), Inches(5.1), Inches(1.3), C_EJ2)
insight_card(s, "4 · DEEPER ES SIEMPRE PEOR",
    ["Última en las 3 columnas de opt.",
     "Profundidad sin batch-norm penaliza la propagación de gradientes."],
    Inches(7.9), Inches(5.75), Inches(5.1), Inches(1.0), C_EJ2)
takeaway_box(s, "Adam domina las 4 archs; wider es especialista de Adam; shallow es robusta al swap de opt; deeper es siempre peor.", C_EJ2)

# ----- Ej2: family comparison -----
s = new_slide("EJ2 · Comparativa head-to-head de optimizadores", C_EJ2)
title_block(s, "Mejor configuración de cada familia", "Lo que se gana o se pierde con cada optimizador")
add_image(s, ROOT/"Notas/ejercicio 2/Segunda tanda de experimentos/Cross_LR_Opt_Arch/family_comparison_bars.png",
          Inches(0.4), Inches(1.55), width=Inches(7.3))
add_table(s,
    ["Familia", "best(arch@LR)", "val_acc", "val_loss CE", "best_epoch", "gap"],
    [
        ["SGD",      "shallow @ 1e-2", "0.9509 ± 0.0053", "0.194", "18.6", "0.163"],
        ["Momentum", "shallow @ 1e-2", "0.9543 ± 0.0062", "0.234",  "4.7", "0.215"],
        ["Adam",     "wider @ 1e-3",   "0.9583 ± 0.0036", "0.170",  "3.4", "0.155"],
    ],
    Inches(7.85), Inches(1.55), Inches(5.2), Inches(2.0))
insight_card(s, "ADAM GANA EN LAS 3 MÉTRICAS QUE IMPORTAN",
    ["Accuracy · calibración (CE) · sobreajuste (gap) — Y además es más rápido que SGD (~4× en wall-clock).",
     "Única ventaja real de Mom/SGD: robustez al HP — útil si NO se hizo HP search. Acá ese costo ya está pagado."],
    Inches(7.85), Inches(3.75), Inches(5.2), Inches(2.65), C_EJ2)
takeaway_box(s, "Adam gana en accuracy, calibración, sobreajuste y velocidad. El HP search ya pagó el costo de la fragilidad de Adam.", C_EJ2)

# ----- Ej2: stage 2b -----
s = new_slide("EJ2 · Stage 2b · estrella batch", C_EJ2)
title_block(s, "Estrella batch sobre el centro", "Confirmar que batch=64 es el óptimo (no un máximo entre puntos espaciados)")
hp_box(s, [
    "Centro: shallow + Adam + LR=1e-3",
    "Variar SÓLO batch: 16, 32, 64, 128, 256",
    "3 seeds × k=5 = 15 corridas / batch",
    "Resto fijo en el centro",
    "ES patience = 20",
], top=Inches(1.55), height=Inches(2.5))
add_image(s, ROOT/"Notas/ejercicio 2/Segunda tanda de experimentos/Cross_LR_Opt_Arch/stage2b_val_acc_vs_batch.png",
          Inches(4.9), Inches(1.55), width=Inches(8.0))
add_table(s,
    ["batch", "val_acc (15 corridas)", "macro_F1", "val_loss CE", "best_epoch"],
    [
        ["16",  "0.9540 ± 0.0036", "0.8477", "0.1785", "2.3"],
        ["32",  "0.9568 ± 0.0054", "0.8510", "0.1700", "3.9"],
        ["64",  "0.9572 ± 0.0041 ★", "0.8521", "0.1701", "5.7"],
        ["128", "0.9556 ± 0.0045", "0.8502", "0.1742", "7.8"],
        ["256", "0.9548 ± 0.0045", "0.8493", "0.1771", "11.2"],
    ],
    Inches(0.4), Inches(4.2), Inches(12.6), Inches(2.2))
takeaway_box(s, "Curva unimodal con pico claro en batch=64 — el óptimo del pre-experimento queda confirmado al detalle.", C_EJ2)

# ----- Ej2: tiebreaker -----
s = new_slide("EJ2 · Arch tiebreaker", C_EJ2)
title_block(s, "Wider vs Shallow — desempate estadístico", "Top-2 del grid empatadas con Δ=0.0011 · subimos a 15 seeds para distinguir")
hp_box(s, [
    "Tiebreaker: shallow vs wider",
    "Opt: Adam · LR: {5e-4, 1e-3}",
    "12 seeds NUEVOS + 3 previos = 15",
    "k=5 → 75 corridas/celda",
    "SEM ≈ 0.0006 (resuelve Δ ≥ 0.0014)",
], top=Inches(1.55), height=Inches(2.5))
add_image(s, ROOT/"Notas/ejercicio 2/Segunda tanda de experimentos/Arch_tiebreaker/tiebreaker_val_acc.png",
          Inches(4.9), Inches(1.55), width=Inches(8.0))
stat_card(s, "z = 0.65", "shallow vs wider",   Inches(0.4), Inches(4.2), Inches(4.0), Inches(1.3), C_EJ2)
stat_card(s, "0.9576",   "shallow · 75 corridas", Inches(4.55), Inches(4.2), Inches(4.0), Inches(1.3), C_EJ2)
stat_card(s, "0.9581",   "wider · 75 corridas",   Inches(8.7), Inches(4.2), Inches(4.0), Inches(1.3), C_EJ2)
insight_card(s, "RESOLUCIÓN POR OCCAM",
    ["Empate estadístico (necesitábamos |z| > 1.96). Entre dos archs indistinguibles, se elige la de MENOS parámetros:",
     "shallow ≈ 101k params  ·  wider ≈ 235k params  →  GANADOR: shallow."],
    Inches(0.4), Inches(5.65), Inches(12.6), Inches(0.85), C_EJ2)
takeaway_box(s, "z = 0.65 → empate estadístico. Occam decide: shallow (menos params, menos riesgo de overfit).", C_EJ2)

# ----- Ej2: configuración óptima -----
s = new_slide("EJ2 · Configuración óptima", C_EJ2)
title_block(s, "Configuración óptima del Ej2", "Decisión congelada al cerrar la 2ª tanda")
# big box with the winning config
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.55), Inches(12.5), Inches(1.7))
box.adjustments[0] = 0.06
box.fill.solid(); box.fill.fore_color.rgb = C_EJ2; box.line.fill.background()
tf = box.text_frame
tf.margin_top = Inches(0.2); tf.margin_bottom = Inches(0.15)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.clear()
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "arch_shallow  [784, 128, 10]   +   Adam  (β1=0.9, β2=0.999, ε=1e-8)"
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r = p2.add_run(); r.text = "LR = 1e-3      ·      batch = 64"
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE

insight_card(s, "CADENA DE DECISIONES",
    ["1. Grid LR×Opt×Arch → Adam domina las 4 archs.",
     "2. Dentro de Adam, top-2 = wider y shallow.",
     "3. Tiebreaker 75 corridas → z=0.65 → EMPATE estadístico.",
     "4. OCCAM → arch_shallow (101k vs 235k params).",
     "5. LR óptimo de Adam = 1e-3 en las 4 archs.",
     "6. batch=64 confirmado por estrella batch."],
    Inches(0.4), Inches(3.4), Inches(7.2), Inches(3.0), C_EJ2)
insight_card(s, "MÉTRICAS DEL CV INTERNO (15 corridas)",
    ["val_acc       = 0.9572 ± 0.0041",
     "macro_F1      = 0.8521 ± 0.0067",
     "val_loss CE   = 0.1701",
     "best_epoch    = 5.7 ± 1.6",
     "",
     "Promedio sobre 3 seeds × 5 folds estratificados."],
    Inches(7.7), Inches(3.4), Inches(5.3), Inches(3.0), C_EJ2)

# ----- Ej2: convergencia óptima -----
s = new_slide("EJ2 · Convergencia de la config óptima", C_EJ2)
title_block(s, "Convergencia del óptimo · CV interno", "15 epoch_histories agregados (3 seeds × 5 folds)")
add_image(s, ROOT/"Notas/ejercicio 2/Segunda tanda de experimentos/Cross_LR_Opt_Arch/optimal_convergence.png",
          Inches(0.4), Inches(1.55), width=Inches(8.2))
insight_card(s, "ES BIEN CALIBRADO",
    ["best_epoch = 5.7 ± 1.6",
     "stop_epoch = 25.7 ± 1.6  ≈  best + patience(20)",
     "0/15 corridas tocaron max_epochs=40."],
    Inches(8.8), Inches(1.55), Inches(4.2), Inches(1.7), C_EJ2)
insight_card(s, "TRAIN MEMORIZA",
    ["train_loss: 0.17 → 0.005 en ~10 ép.",
     "train_acc → 0.9979 (capacidad de sobra)."],
    Inches(8.8), Inches(3.3), Inches(4.2), Inches(1.5), C_EJ2)
insight_card(s, "VAL TIENE TECHO",
    ["val_loss mínimo ≈ 0.170 en ép. 5.7.",
     "Después sube ligeramente.",
     "→ techo del dataset, no del modelo."],
    Inches(8.8), Inches(4.85), Inches(4.2), Inches(1.5), C_EJ2)
takeaway_box(s, "val_acc = 0.957 NO es techo de Adam ni del LR. Es techo del DATASET sin regularización — motiva el Ej3.", C_EJ2)

# ----- Ej2: generalización CV interno -----
s = new_slide("EJ2 · Generalización interna · CV", C_EJ2)
title_block(s, "Generalización interna (CV sobre digits.csv)", "Métricas del fold de validación · media sobre 15 corridas")
add_table(s,
    ["Métrica", "Train (CV)", "Val (CV)", "Gap val − train"],
    [
        ["accuracy",         "0.9979 ± 0.0023",  "0.9572 ± 0.0041", "−0.041"],
        ["macro_precision",  "(no almacenada)",  "0.8546 ± 0.0062", "—"],
        ["macro_recall",     "(no almacenada)",  "0.8502 ± 0.0078", "—"],
        ["macro_F1",         "(no almacenada)",  "0.8521 ± 0.0067", "—"],
        ["CE loss",          "0.0180 ± 0.0104",  "0.1701 ± 0.0170", "+0.152"],
        ["best_epoch",       "—",                "5.7 ± 1.6",       "—"],
    ],
    Inches(0.4), Inches(1.55), Inches(12.6), Inches(2.8))
insight_card(s, "DIAGNÓSTICO",
    ["• train_acc = 0.998 → memoria casi perfecta → CAPACIDAD DE SOBRA.",
     "• Brecha acc ↔ F1 (~10pp) es ESTRUCTURAL: clase 5 con 271 ejemplos arrastra macro_F1.",
     "• Gap train → val: +0.041 acc, +0.152 CE → OVERFIT MODERADO esperable para ~101k params sobre 10k samples sin reg."],
    Inches(0.4), Inches(4.5), Inches(12.6), Inches(1.9), C_EJ2)
takeaway_box(s, "El modelo memoriza train pero el techo en val es del dataset, no del optimizador. F1 macro bajo = clase 5 desbalanceada.", C_EJ2)

# ----- Ej2: generalización externa -----
s = new_slide("EJ2 · Generalización externa · digits_test", C_EJ2)
title_block(s, "Final eval sobre digits_test.csv", "3 seeds · config congelada · evaluado UNA sola vez")
add_image(s, ROOT/"Notas/ejercicio 2/Segunda tanda de experimentos/Cross_LR_Opt_Arch/optimal_test_confusion_matrix.png",
          Inches(0.4), Inches(1.55), width=Inches(6.8))
stat_card(s, "−10 pp", "Δ accuracy val → test", Inches(7.4), Inches(1.55), Inches(5.6), Inches(1.4), C_EJ2)
stat_card(s, "0.853", "test accuracy",          Inches(7.4), Inches(3.1),  Inches(2.7), Inches(1.3), C_EJ2)
stat_card(s, "0.806", "test macro_F1",          Inches(10.25),Inches(3.1), Inches(2.75),Inches(1.3), C_EJ2)
insight_card(s, "CAUSA ESTRUCTURAL — CLASE 8 AUSENTE",
    ["digits.csv tiene 0 ejemplos de clase 8.",
     "digits_test.csv tiene 243 (9.7%) → IMPOSIBLES de acertar.",
     "Excluyendo la clase 8: test_acc = 0.9448 ± 0.003",
     "→ gap residual REAL = +1.2 pp."],
    Inches(7.4), Inches(4.55), Inches(5.6), Inches(1.85), C_EJ2)
takeaway_box(s, "La caída de 10 pp es estructural del DATASET (clase 8 ausente), no del modelo. La intervención correcta: MÁS DATOS — Ej3.", C_EJ2)

# ===============================================================
# EJ3 — DIVISOR
# ===============================================================
PAGE += 1
s = prs.slides.add_slide(BLANK)
section_divider(s, "Ejercicio 3", "Hacia ≥ 98% sobre digits_test.csv · más datos + regularización", C_EJ3, PAGE)

# ----- Ej3: plan -----
s = new_slide("EJ3 · Plan y motivación", C_EJ3)
title_block(s, "Plan y motivación", "Por qué regularización y no “otra arquitectura”")
insight_card(s, "PALANCA 1 — MÁS DATOS",
    ["Sumar more_digits.csv al training:",
     "• 585 ejemplos de la clase 8 (antes ausente)",
     "• +542 ejemplos de la clase 5 (271 → 813)",
     "• Total: 12.449 → 28.190 muestras",
     "",
     "Hipótesis: test_acc sube a 0.93-0.95."],
    Inches(0.4), Inches(1.55), Inches(6.2), Inches(2.6), C_EJ3, body_size=12)
insight_card(s, "PALANCA 2 — REGULARIZACIÓN (de la clase)",
    ["Grid 4 × 4 sobre el baseline de palanca 1:",
     "• L2 (weight decay): {0, 1e-5, 1e-4, 1e-3}",
     "  → slides 20-25 de regularización (fórmula exacta)",
     "• σ (gaussian noise): {0, 0.03, 0.1, 0.2}",
     "  → slide 18 de augmentation"],
    Inches(6.8), Inches(1.55), Inches(6.2), Inches(2.6), C_EJ3, body_size=12)
insight_card(s, "QUÉ DEJAMOS AFUERA Y POR QUÉ",
    ["• Dropout: mencionado en la clase como “existen otros más”, no profundizado → fuera del grid.",
     "• LR schedule: no aparece en la clase → fuera del grid.",
     "→ El criterio es estricto: sólo técnicas profundizadas en las clases del curso."],
    Inches(0.4), Inches(4.25), Inches(12.6), Inches(1.5), C_EJ3, body_size=12)
stat_card(s, "240", "corridas CV totales", Inches(0.4), Inches(5.9), Inches(4.0), Inches(0.55), C_EJ3)
stat_card(s, "16",  "combos del grid",     Inches(4.55), Inches(5.9), Inches(4.0), Inches(0.55), C_EJ3)
stat_card(s, "3 × 5", "seeds × folds",      Inches(8.7), Inches(5.9), Inches(4.0), Inches(0.55), C_EJ3)
takeaway_box(s, "Dos palancas, ambas en la clase: sumar more_digits.csv y grid L2 × σ. Dropout y LR schedule quedan fuera por criterio.", C_EJ3)

# ----- Ej3: baseline +more_digits convergencia -----
s = new_slide("EJ3 · Paso 1 · Baseline + more_digits", C_EJ3)
title_block(s, "Paso 1 · Baseline con more_digits.csv (sin reg)", "Aislar el efecto de “más datos” manteniendo el resto del Ej2")
hp_box(s, [
    "Arch: shallow [784, 128, 10]",
    "Adam @ 1e-3 · batch = 64",
    "Datasets: digits.csv + more_digits.csv",
    "  N = 12.449 + 15.741 = 28.190",
    "k=5 estratificado · 3 seeds (42, 7, 13)",
    "max_epochs = 50 · ES patience = 20",
    "Regularización: todo en 0",
], top=Inches(1.55), height=Inches(3.0))
add_image(s, ROOT/"ejercicio3/analisis/baseline/optimal_convergence.png",
          Inches(4.9), Inches(1.55), width=Inches(8.0))
insight_card(s, "CONVERGENCIA (15 CORRIDAS)",
    ["best_epoch  = 5.7 ± 1.1   ·   stop_epoch  = 25.7 ± 1.1",
     "0/15 corridas tocaron max_epochs=50 → ES bien calibrado.",
     "Mismo patrón que Ej2 — agregar más datos no cambió la velocidad de convergencia."],
    Inches(0.4), Inches(4.65), Inches(12.6), Inches(1.7), C_EJ3, body_size=12)
takeaway_box(s, "Sumar more_digits.csv no cambia la dinámica de entrenamiento — convergencia y ES idénticos al Ej2.", C_EJ3)

# ----- Ej3: baseline test -----
s = new_slide("EJ3 · Paso 1 · Generalización en test", C_EJ3)
title_block(s, "+more_digits.csv → salto en test", "Comparativa Ej2 vs Ej3 baseline")
add_image(s, ROOT/"ejercicio3/analisis/baseline/test_confusion_matrix_baseline.png",
          Inches(0.4), Inches(1.55), width=Inches(6.6))
add_table(s,
    ["Configuración", "Test accuracy", "Test macro_F1"],
    [
        ["Ej2 (sin more_digits, sin reg)",      "0.8529 ± 0.0034", "0.8062 ± 0.0034"],
        ["Ej3 baseline (+more_digits, sin reg)","0.9616 ± 0.0025", "0.9609 ± 0.0026"],
        ["Δ (Ej3 − Ej2)",                       "+0.1087",          "+0.1547"],
    ],
    Inches(7.2), Inches(1.55), Inches(5.8), Inches(2.0))
stat_card(s, "+10.87 pp", "ganancia test_acc", Inches(7.2), Inches(3.75), Inches(5.8), Inches(1.2), C_EJ3)
insight_card(s, "CLASE 8 — AHORA SE APRENDE",
    ["F1 clase 8 = 0.938 (prec=0.968, rec=0.909).",
     "Test acc EXCLUYENDO clase 8 = 0.9672 ± 0.0017.",
     "→ el shift residual entre digits.csv y test es mínimo."],
    Inches(7.2), Inches(5.1), Inches(5.8), Inches(1.3), C_EJ3, body_size=11)
takeaway_box(s, "Sumar more_digits.csv mueve la aguja +10.87 pp — la palanca dominante del Ej3.", C_EJ3)

# ----- Ej3: grid heatmap val_acc -----
s = new_slide("EJ3 · Paso 2 · Grid L2 × σ", C_EJ3)
title_block(s, "Grid de regularización L2 × σ", "16 combos × 3 seeds × 5 folds = 240 corridas CV")
hp_box(s, [
    "Sobre el baseline del paso 1",
    "L2:  {0, 1e-5, 1e-4, 1e-3}",
    "σ:   {0, 0.03, 0.1, 0.2}",
    "(píxeles z-score → σ = fracción del std)",
    "Loss penalizada: CE + (λ/2)‖W‖² (no bias)",
    "Aug: ruido N(0, σ) por minibatch",
    "Resto = baseline paso 1",
], top=Inches(1.55), height=Inches(3.3))
add_image(s, ROOT/"ejercicio3/analisis/grid_reg/val_acc_heatmap.png",
          Inches(4.9), Inches(1.55), width=Inches(8.0))
insight_card(s, "BEST COMBO CV",
    ["L2 = 1e-3 · σ = 0 → val_acc = 0.9750 ± 0.0018 (gap 0.0750).",
     "L2 mueve la aguja MÁS que σ: las 4 filas con L2=1e-3 superan al baseline (L2=0).",
     "σ alto (0.2) no degrada en general — curva no monótona, suele aplanar el efecto de σ medio."],
    Inches(0.4), Inches(5.0), Inches(12.6), Inches(1.5), C_EJ3, body_size=11)
takeaway_box(s, "Best CV: L2=1e-3, σ=0 → val_acc 0.9750 (+0.005 vs baseline sin reg). L2 domina sobre σ en este grid.", C_EJ3)

# ----- Ej3: gap heatmap -----
s = new_slide("EJ3 · Heatmap del gap val − train", C_EJ3)
title_block(s, "Heatmap del sobreajuste (val − train CE)", "Confirma que la regularización REDUCE el gap")
add_image(s, ROOT/"ejercicio3/analisis/grid_reg/gap_heatmap.png", Inches(0.4), Inches(1.55), width=Inches(7.2))
stat_card(s, "0.112", "gap CE · sin reg",       Inches(7.9), Inches(1.55), Inches(5.1), Inches(1.3), C_EJ3)
stat_card(s, "0.075", "gap CE · L2=1e-3",        Inches(7.9), Inches(3.05), Inches(5.1), Inches(1.3), C_EJ3)
stat_card(s, "−33%",  "reducción del gap",       Inches(7.9), Inches(4.55), Inches(5.1), Inches(1.3), C_EJ3)
insight_card(s, "QUÉ DICE EL HEATMAP",
    ["L2=1e-3 cierra ~30% del sobreajuste manteniendo la accuracy.",
     "σ por sí solo (L2=0) mueve poco el gap."],
    Inches(7.9), Inches(6.0), Inches(5.1), Inches(0.45), C_EJ3, body_size=10)
takeaway_box(s, "L2=1e-3 cierra ~30% del gap sin perder accuracy en CV. σ sola no cambia el gap.", C_EJ3)

# ----- Ej3: best reg test -----
s = new_slide("EJ3 · Best combo en test", C_EJ3)
title_block(s, "Best combo evaluado en digits_test.csv", "L2=1e-3, σ=0 · 3 seeds · vs baseline sin reg")
add_image(s, ROOT/"ejercicio3/analisis/best_reg/test_confusion_matrix_best_reg.png",
          Inches(0.4), Inches(1.55), width=Inches(6.6))
add_table(s,
    ["Métrica", "Best reg (test)"],
    [
        ["accuracy",        "0.9601 ± 0.0030"],
        ["macro_precision", "0.9605 ± 0.0029"],
        ["macro_recall",    "0.9591 ± 0.0030"],
        ["macro_F1",        "0.9594 ± 0.0030"],
    ],
    Inches(7.2), Inches(1.55), Inches(5.8), Inches(1.9))
insight_card(s, "OBSERVACIÓN HONESTA",
    ["best_reg en TEST queda −0.0015 PEOR que el baseline (+more_digits, sin reg):",
     "   baseline +more_digits → 0.9616",
     "   best_reg              → 0.9601",
     "",
     "Hipótesis: el gap CV → test acá NO es por sobreajuste de train, es por SHIFT de distribución",
     "entre digits.csv y digits_test.csv. La regularización combate memoria, no shift."],
    Inches(7.2), Inches(3.65), Inches(5.8), Inches(2.75), C_EJ3, body_size=11)
takeaway_box(s, "Best reg en TEST queda un pelo peor que el baseline: el gap residual era SHIFT, no overfit — y la reg no lo combate.", C_EJ3)

# ----- Ej3: comparativa final -----
s = new_slide("EJ3 · Comparativa final", C_EJ3)
title_block(s, "Comparativa final · Ej2 → +more_digits → +reg", "Cuánto aporta cada palanca al objetivo de 98%")
add_table(s,
    ["Configuración", "Test accuracy", "Test macro_F1"],
    [
        ["Ej2 (sin more_digits, sin reg)",              "0.8529 ± 0.0034", "0.8062 ± 0.0034"],
        ["Ej3 baseline (+more_digits, sin reg)",         "0.9616 ± 0.0025", "0.9609 ± 0.0026"],
        ["Ej3 best_reg (+more_digits + L2=1e-3 + σ=0)",  "0.9601 ± 0.0030", "0.9594 ± 0.0030"],
    ],
    Inches(0.4), Inches(1.55), Inches(12.6), Inches(2.0))
stat_card(s, "+10.87 pp", "+more_digits",        Inches(0.4),  Inches(3.7), Inches(4.0), Inches(1.3), C_EJ3)
stat_card(s, "−0.15 pp",  "+L2 sobre +more_dig", Inches(4.55), Inches(3.7), Inches(4.0), Inches(1.3), C_EJ3)
stat_card(s, "0.0199",    "brecha al 98% pedido",Inches(8.7),  Inches(3.7), Inches(4.3), Inches(1.3), C_EJ3)
insight_card(s, "LECTURA DEFENSIVA",
    ["• Palanca dominante: MÁS DATOS (la clase 8 ausente era el agujero estructural).",
     "• Adam ya tenía el gap más bajo de los 3 opts (lo vimos en LR×OPT del Ej2) → reg movía poco la aguja sobre Adam.",
     "• Hipótesis del residuo: el shift entre digits.csv y digits_test.csv es GEOMÉTRICO (rotaciones, traslaciones",
     "  entre escritores). σ gaussiano simula ruido isotrópico (slide 18) — geometría requeriría augmentación",
     "  geométrica, que la clase no profundiza."],
    Inches(0.4), Inches(5.2), Inches(12.6), Inches(1.25), C_EJ3, body_size=11)
takeaway_box(s, "+10.87 pp por más datos · −0.15 pp por reg. La palanca dominante fue SUMAR DATOS — no llegamos al 98%, brecha = 1.99 pp.", C_EJ3)

# ===============================================================
# CIERRE
# ===============================================================
s = new_slide("CIERRE", C_GEN)
title_block(s, "Qué aprendimos · qué quedó pendiente", "Una lectura por ejercicio")
insight_card(s, "EJ1 — La métrica importa",
    ["El no-lineal tiene 2× mejor MSE",
     "que el lineal pero quedan empatados",
     "en F1, ambos rozando el baseline.",
     "",
     "La loss que se minimiza (MSE) NO",
     "responde la misma pregunta que la",
     "métrica de decisión (F1) — regla 4."],
    Inches(0.4), Inches(1.55), Inches(4.2), Inches(4.6), C_EJ1, body_size=12)
insight_card(s, "EJ2 — Los HP interactúan",
    ["El “techo” de Adam no era del opt:",
     "era del producto LR×batch.",
     "El ranking de archs depende del LR —",
     "un sweep one-at-a-time mentía.",
     "",
     "Decisión final defendida con criterio",
     "explícito: empate estadístico → Occam."],
    Inches(4.7), Inches(1.55), Inches(4.2), Inches(4.6), C_EJ2, body_size=12)
insight_card(s, "EJ3 — Datos > regularización",
    ["+more_digits subió +10.87 pp;",
     "la regularización L2+σ movió −0.15.",
     "",
     "El residuo es SHIFT geométrico que",
     "L2 y σ gaussiano no resuelven.",
     "Augmentación geométrica (slide 18 de",
     "regularización) no se profundizó."],
    Inches(9.0), Inches(1.55), Inches(3.9), Inches(4.6), C_EJ3, body_size=12)
takeaway_box(s, "Reglas que sostuvimos: loss + Acc/Prec/Rec/F1 siempre · ejes de promedio explícitos · decisiones ancladas en las clases.", C_GEN)

# ---------- SAVE ----------
prs.save(str(OUT))
print(f"OK: {OUT}  ·  slides = {len(prs.slides)}")
